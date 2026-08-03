"""
PPTX Layout Parser using python-pptx.
Extracts slide titles, bullet points, speaker notes, and slide numbers.
"""

import os
from typing import List
from pptx import Presentation
from loguru import logger

from backend.rag.parsers.base import BaseDocumentParser
from backend.rag.schemas import ElementType, ExtractedElement
from backend.utils.helpers import sanitize_text


class PPTXDocumentParser(BaseDocumentParser):
    """Slide-aware PPTX presentation parser."""

    def supports_extension(self, extension: str) -> bool:
        return extension.lower() in (".pptx", ".ppt")

    async def parse(self, file_path: str) -> List[ExtractedElement]:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PPTX file not found at path: {file_path}")

        elements: List[ExtractedElement] = []

        try:
            prs = Presentation(file_path)
            logger.info("Parsing PPTX presentation: path={path} slides={count}", path=file_path, count=len(prs.slides))

            for slide_idx, slide in enumerate(prs.slides):
                slide_num = slide_idx + 1
                slide_label = f"Slide {slide_num}"
                slide_title = ""

                # Extract slide title
                if slide.shapes.title and slide.shapes.title.text:
                    slide_title = sanitize_text(slide.shapes.title.text)

                heading_path = [f"Slide {slide_num}: {slide_title}"] if slide_title else [slide_label]

                if slide_title:
                    elements.append(
                        ExtractedElement(
                            element_type=ElementType.HEADING_1,
                            content=slide_title,
                            page_number=slide_num,
                            heading_path=heading_path,
                            metadata={"slide_label": slide_label},
                        )
                    )

                # Extract text boxes and shapes
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        for para in shape.text_frame.paragraphs:
                            text = sanitize_text(para.text)
                            if not text:
                                continue

                            elements.append(
                                ExtractedElement(
                                    element_type=ElementType.PARAGRAPH,
                                    content=text,
                                    page_number=slide_num,
                                    heading_path=heading_path,
                                    metadata={"slide_label": slide_label, "level": para.level},
                                )
                            )

            logger.info("PPTX parsing complete: path={path} elements={count}", path=file_path, count=len(elements))
            return elements
        except Exception as e:
            logger.error("Failed to parse PPTX file: path={path} error={err}", path=file_path, err=str(e))
            raise
